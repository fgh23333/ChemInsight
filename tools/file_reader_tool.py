import io
import os
from google.genai import types
from google.adk.agents import callback_context
from google.adk.models.llm_request import LlmRequest

# --- 导入所有需要的文件处理库 (如果库未安装，会给出明确错误提示) ---
try:
    import docx
except ImportError:
    class docx:
        @staticmethod
        def Document(stream):
            raise ImportError("请运行 'pip install python-docx' 来安装处理 .docx 文件的库。")
try:
    import openpyxl
except ImportError:
    class openpyxl:
        @staticmethod
        def load_workbook(stream):
            raise ImportError("请运行 'pip install openpyxl' 来安装处理 .xlsx 文件的库。")
try:
    import pptx
except ImportError:
    class pptx:
        @staticmethod
        def Presentation(stream):
            raise ImportError("请运行 'pip install python-pptx' 来安装处理 .pptx 文件的库。")

# --- 核心修改：将 application/pdf 添加到原生支持的类型中 ---
# 我们的回调将跳过这些类型，让模型直接处理
NATIVELY_SUPPORTED_MIME_PREFIXES = ("image", "video", "application/pdf")

def read_files_as_text_callback(
    callback_context: callback_context, llm_request: LlmRequest
):
    """
    一个“模型调用前”回调函数，用于智能预处理上传的各类文件。

    - **新增**：PDF 文件 (application/pdf) 将被跳过，由模型直接处理。
    - **解决了因MIME类型为空导致Markdown文件被过滤的问题。**
    - **解决了因访问不存在的.uri属性导致的AttributeError。**
    - 智能地结合MIME类型和文件名后缀来判断文件类型。
    - 为所有提取出的文本提供了清晰的上下文，以获得更准确的模型响应。
    """
    if not llm_request.contents:
        return

    for content in llm_request.contents:
        if content.role != "user":
            continue

        new_parts = []
        for part in content.parts:
            # 检查这是否是一个包含文件数据的 part
            if not (part.inline_data and part.inline_data.data):
                new_parts.append(part)
                continue

            mime_type = part.inline_data.mime_type or ""

            # 如果文件类型是原生支持的（图片、视频、PDF），则直接跳过，保留原始part
            if any(mime_type.startswith(prefix) for prefix in NATIVELY_SUPPORTED_MIME_PREFIXES):
                new_parts.append(part)
                continue

            # --- 进入我们的文件处理逻辑 ---
            text_content = None
            file_name = part.inline_data.display_name or "uploaded_file"
            
            try:
                # 结合MIME类型和文件名后缀进行判断
                file_ext = os.path.splitext(file_name)[1].lower()

                # 1. 处理 Word (.docx)
                if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or file_ext == ".docx":
                    document = docx.Document(io.BytesIO(part.inline_data.data))
                    text_content = "\n".join([p.text for p in document.paragraphs])
                # 2. 处理 Excel (.xlsx)
                elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" or file_ext == ".xlsx":
                    workbook = openpyxl.load_workbook(io.BytesIO(part.inline_data.data))
                    full_text = [f"--- 工作表: {name} ---\n" + "\n".join("\t".join(str(cell.value) if cell.value is not None else "" for cell in row) for row in workbook[name].iter_rows()) for name in workbook.sheetnames]
                    text_content = "\n\n".join(full_text)
                # 3. 处理 PowerPoint (.pptx)
                elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or file_ext == ".pptx":
                    presentation = pptx.Presentation(io.BytesIO(part.inline_data.data))
                    full_text = [f"--- 幻灯片 {i+1} ---\n" + "\n".join(shape.text_frame.text for shape in slide.shapes if hasattr(shape, "text_frame") and shape.text_frame) for i, slide in enumerate(presentation.slides)]
                    text_content = "\n\n".join(full_text)
                # 4. 如果不是Office文档，则一律尝试作为文本文件解码
                else:
                    text_content = part.inline_data.data.decode("utf-8")

                if text_content is not None:
                    formatted_text = (
                        f"用户上传了文件 '{file_name}' (MIME类型: '{mime_type}'), 其提取出的文本内容如下:\n--- 文件开始 ---\n"
                        f"{text_content}\n--- 文件结束 ---"
                    )
                    new_parts.append(types.Part(text=formatted_text))
                else:
                    raise ValueError("未能从文件中提取任何文本内容")

            except Exception as e:
                error_message = f"无法自动读取文件 '{file_name}' 的内容 (MIME类型: {mime_type})，错误: {e}。请基于文件名进行推断，或告知用户无法处理此文件。"
                new_parts.append(part)
                new_parts.append(types.Part(text=error_message))
        
        content.parts = new_parts